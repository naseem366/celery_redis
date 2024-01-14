from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
from PIL import Image
import os

def add_slide(prs, slide_data):
    template_id = json_data["data"]["templateId"]
    print(template_id)
    if template_id == 102:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
    
        # Title
        title = slide.shapes.title
        title.text = slide_data["title"]
    
        # Content with bullet points on the right side
        content_box = slide.placeholders[1]
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_box.left = Inches(0.5)
        content_box.top = Inches(1)
        content_box.width = Inches(6)
    
        # # Set font properties for content text
        # font_size = Pt(18)  # Adjust the font size as needed
        # font_name = 'Arial'  # Change to your desired font family
        # for paragraph in content_box.text_frame.paragraphs:
        #     for run in paragraph.runs:
        #     run.font.size = font_size
        # run.font.name = font_name
        
        content_lines = slide_data["content"].split('\n')
        for line in content_lines:
            p = content_frame.add_paragraph()
            p.text = line
            p.space_after = Pt(14)  # Adjust the spacing as needed
            p.level = 0  # Level 0 indicates a bullet point
         
        # Image on the left side
        image_path_or_url = slide_data["image_url"]
        if image_path_or_url.startswith("http"):
            # If the image is from a URL
            image_response = requests.get(image_path_or_url)
            try:
                image = Image.open(BytesIO(image_response.content))
                img_width, img_height = image.size
                
                # Adjust the left position as needed
                image_width = Inches(4)
                image_height = Inches(4)
                image_left = Inches(6)
                image_top = Inches(3)
                img = slide.shapes.add_picture(BytesIO(image_response.content),image_left,image_top,width=image_width,height=image_height)
            except Exception as e:
                print(f"Error adding image from URL: {e}")
        elif os.path.exists(image_path_or_url):
            # If the image is a local file
            try:
                image = Image.open(image_path_or_url)
                img_width, img_height = image.size
                left_inch = Inches(4.5)  # Adjust the left position as needed
                top_inch = Inches(1)
                img = slide.shapes.add_picture(image_path_or_url, Inches(0.5), Inches(4.3), width=Inches(3))
            except Exception as e:
                print(f"Error adding image from path: {e}")
        else:
            print(f"Image not found: {image_path_or_url}")  

############# with TemplateId ==103 
    
    elif template_id == 103:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
    
        # Title
        title = slide.shapes.title
        title.text = slide_data["title"]
    
        # Content with bullet points on the right side
        content_box = slide.placeholders[1]
        content_box.text = slide_data["content"]
        #content_box = slide.placeholders[1]
        #content_frame = content_box.text_frame
        #content_frame.word_wrap = True
        content_box.left = Inches(4)
        content_box.top = Inches(2)
        content_box.width = Inches(6)
    
        # Set font properties for content text
        font_size = Pt(16)  # Adjust the font size as needed
        font_name = 'Arial'  # Change to your desired font family
        font_color = RGBColor(255, 0, 0)  # Red color, adjust as needed
        font_bold = False  # Set to True for bold, False for unbold
        
        for paragraph in content_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.name = font_name
                run.font.color.rgb = font_color
                run.font.bold = font_bold
        
        # content_lines = slide_data["content"].split('\n')
        # for line in content_lines:
        #     p = content_frame.add_paragraph()
        #     p.text = line
        #     p.space_after = Pt(14)  # Adjust the spacing as needed
        #     p.level = 0  # Level 0 indicates a bullet point
         
        # Image on the left side
        image_path_or_url = slide_data["image_url"]
        if image_path_or_url.startswith("http"):
            # If the image is from a URL
            image_response = requests.get(image_path_or_url)
            try:
                image = Image.open(BytesIO(image_response.content))
                img_width, img_height = image.size
                
                # Adjust the left position as needed
                image_width = Inches(3.5)
                image_height = Inches(3.5)
                image_left = Inches(0.5)
                image_top = Inches(2.5)
                img = slide.shapes.add_picture(BytesIO(image_response.content),image_left,image_top,width=image_width,height=image_height)
            except Exception as e:
                print(f"Error adding image from URL: {e}")
        elif os.path.exists(image_path_or_url):
            # If the image is a local file
            try:
                image = Image.open(image_path_or_url)
                img_width, img_height = image.size
                left_inch = Inches(4.5)  # Adjust the left position as needed
                top_inch = Inches(1)
                img = slide.shapes.add_picture(image_path_or_url, Inches(0.5), Inches(4.3), width=Inches(3))
            except Exception as e:
                print(f"Error adding image from path: {e}")
        else:
            print(f"Image not found: {image_path_or_url}")
    else:
        print("Template Id not exists")

if __name__ == "__main__":
    # JSON data
    json_data = {
    "data": {
        "templateId": 103,
        "fileName" : "DemoPPT",
        "slides": [
            {
                "slide_number": "1",
                "title": "Slide 1 Title",
                "content": "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Etiam tristique est id ipsum malesuada, ac ultrices nulla interdum. Donec cursus erat id eros efficitur, a vulputate libero tempus. Donec nec leo nec risus tempor placerat. Morbi tincidunt condimentum ut. ",
                "image_suggestion": "health_and_wellness_product_image_1",
                "image_url": "https://picsum.photos/200/300"
            },
            {
                "slide_number": "2",
                "title": "Slide 2 Title",
                "content": "Donec eu sem augue. Maecenas rhoncus orci id felis pharetra accumsan. Etiam mattis, velit quis ultrices scelerisque, eros urna euismod magna, sed tincidunt orci diam id metus. Fusce interdum consectetur sagittis. Nullam at aliquet diam. Fusce dapibus a metus vulputate fringilla. ",
                "image_suggestion": "health_and_wellness_product_image_2",
                "image_url": "https://picsum.photos/200/300"
            }
        ]
    }
}

    # Create PowerPoint presentation
    presentation = Presentation()

    # Add slides to the presentation
    for slide_data in json_data["data"]["slides"]:
        add_slide(presentation, slide_data)

    # Save the presentation
    presentation.save("Hippa_Presentation.pptx")